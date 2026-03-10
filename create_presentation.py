"""
Generate UW-style presentation for CFRM 522 Strategy Project
Fixed image sizing and alignment
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# UW Colors
UW_PURPLE = RGBColor(0x4B, 0x2E, 0x83)
UW_GOLD = RGBColor(0xB7, 0xA5, 0x7A)
UW_LIGHT_GOLD = RGBColor(0xE8, 0xE3, 0xD3)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x2D, 0x2D, 0x2D)
LIGHT_GRAY = RGBColor(0x66, 0x66, 0x66)

# Slide dimensions
SLIDE_WIDTH = 13.333
SLIDE_HEIGHT = 7.5
HEADER_HEIGHT = 0.9
CONTENT_TOP = 1.1
CONTENT_BOTTOM = 6.8
MARGIN = 0.4

def set_background(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def add_text(slide, text, left, top, width, height, size=18, bold=False, color=DARK_GRAY, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Arial"
    p.alignment = align
    return box

def add_bullets(slide, items, left, top, width, height, size=16, color=DARK_GRAY):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Arial"
        p.space_before = Pt(6)
        p.space_after = Pt(6)
    return box

def add_header(slide, title, prs):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(HEADER_HEIGHT))
    bar.fill.solid()
    bar.fill.fore_color.rgb = UW_PURPLE
    bar.line.fill.background()

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(HEADER_HEIGHT), prs.slide_width, Inches(0.03))
    accent.fill.solid()
    accent.fill.fore_color.rgb = UW_GOLD
    accent.line.fill.background()

    add_text(slide, title, Inches(MARGIN), Inches(0.22), Inches(12), Inches(0.5), size=26, bold=True, color=WHITE)

def add_image_fit(slide, path, left, top, max_width, max_height):
    """Add image with proper fitting within bounds."""
    if not os.path.exists(path):
        return
    # Add with height constraint to prevent overflow
    pic = slide.shapes.add_picture(path, left, top, height=max_height)
    # Check if width exceeds max, if so resize by width
    if pic.width > max_width:
        ratio = max_width / pic.width
        pic.width = max_width
        pic.height = int(pic.height * ratio)

def add_caption(slide, text):
    add_text(slide, text, Inches(MARGIN), Inches(CONTENT_BOTTOM), Inches(SLIDE_WIDTH - 2*MARGIN),
             Inches(0.4), size=12, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

def add_section_title(slide, text, left, top):
    add_text(slide, text, left, top, Inches(5), Inches(0.35), size=17, bold=True, color=UW_PURPLE)

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH)
    prs.slide_height = Inches(SLIDE_HEIGHT)

    # ══════════════════════════════════════════════════════════════
    # SLIDE 1: Title
    # ══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, UW_PURPLE)

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(3.0), Inches(11.7), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = UW_GOLD
    line.line.fill.background()

    add_text(slide, "Funding Rate Contrarian Strategy", Inches(0.8), Inches(1.8), Inches(11), Inches(1), size=46, bold=True, color=WHITE)
    add_text(slide, "BTC-USDT Perpetual Futures | Binance", Inches(0.8), Inches(3.3), Inches(11), Inches(0.5), size=22, color=UW_LIGHT_GOLD)
    add_text(slide, "Adeline Wen", Inches(0.8), Inches(5.4), Inches(5), Inches(0.4), size=20, color=WHITE)
    add_text(slide, "CFRM 522 · Winter 2026", Inches(0.8), Inches(5.9), Inches(5), Inches(0.4), size=16, color=UW_LIGHT_GOLD)

    # ══════════════════════════════════════════════════════════════
    # SLIDE 2: Problem & Hypothesis
    # ══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, WHITE)
    add_header(slide, "Problem & Hypothesis", prs)

    add_section_title(slide, "Economic Mechanism", Inches(MARGIN), Inches(CONTENT_TOP))
    add_bullets(slide, [
        "Perpetual futures: funding paid every 8 hours",
        "Baseline 0.01% ≈ 13% annualized cost",
        "Extreme 0.05%+ → 65%+ annualized",
        "Crowded longs forced to exit → price drops",
        "Brunnermeier & Pedersen (2009): liquidity spiral"
    ], Inches(MARGIN), Inches(1.5), Inches(5.8), Inches(2.6), size=14)

    add_section_title(slide, "Pre-committed Hypotheses", Inches(6.8), Inches(CONTENT_TOP))
    add_bullets(slide, [
        "H1: |z-score| > 1.5 predicts 24h reversal",
        "H2: High OI change amplifies signal",
        "H3: Threshold-Calmar is non-monotonic",
        "H4: WF ratio > 0.5 (OOS positive)"
    ], Inches(6.8), Inches(1.5), Inches(6), Inches(2.2), size=14)

    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(MARGIN), Inches(4.3), Inches(12.5), Inches(0.8))
    box.fill.solid()
    box.fill.fore_color.rgb = UW_LIGHT_GOLD
    box.line.fill.background()
    add_text(slide, "Strategy: Short when z > threshold  |  Long when z < -threshold",
             Inches(0.6), Inches(4.5), Inches(12), Inches(0.4), size=18, bold=True, color=UW_PURPLE, align=PP_ALIGN.CENTER)

    add_section_title(slide, "Objective: Calmar Ratio", Inches(MARGIN), Inches(5.4))
    add_text(slide, "BTC excess kurtosis > 20 — Sharpe assumes normality; Calmar penalizes max drawdown directly",
             Inches(MARGIN), Inches(5.8), Inches(12), Inches(0.4), size=13, color=LIGHT_GRAY)

    # ══════════════════════════════════════════════════════════════
    # SLIDE 3: Funding Distribution
    # ══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, WHITE)
    add_header(slide, "Data: Funding Rate Distribution", prs)
    add_image_fit(slide, "data/fig_funding_dist.png",
                  Inches(0.5), Inches(CONTENT_TOP + 0.1), Inches(12.3), Inches(4.8))
    add_caption(slide, "Heavy tails (kurtosis > 10)  •  Positive skew (longs dominate)  •  Non-normal → justifies Calmar over Sharpe")

    # ══════════════════════════════════════════════════════════════
    # SLIDE 4: ACF & Seasonality
    # ══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, WHITE)
    add_header(slide, "Data: Autocorrelation & Seasonality", prs)
    add_image_fit(slide, "data/fig_funding_acf.png",
                  Inches(0.3), Inches(CONTENT_TOP + 0.1), Inches(6.2), Inches(4.2))
    add_image_fit(slide, "data/fig_funding_seasonal.png",
                  Inches(6.7), Inches(CONTENT_TOP + 0.1), Inches(6.2), Inches(4.2))
    add_text(slide, "ACF decays after lag 3 → mean reversion",
             Inches(0.5), Inches(5.5), Inches(6), Inches(0.3), size=12, color=LIGHT_GRAY)
    add_text(slide, "No weekday effect → sentiment-driven",
             Inches(6.8), Inches(5.5), Inches(6), Inches(0.3), size=12, color=LIGHT_GRAY)

    # ══════════════════════════════════════════════════════════════
    # SLIDE 5: Price & Events
    # ══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, WHITE)
    add_header(slide, "Data: BTC Price & Funding Rate (2020-2024)", prs)
    add_image_fit(slide, "data/fig_events.png",
                  Inches(0.3), Inches(CONTENT_TOP), Inches(12.7), Inches(5.0))
    add_caption(slide, "Key events: March 2020 crash  •  May 2021 selloff  •  LUNA (May 2022)  •  FTX (Nov 2022)")

    # ══════════════════════════════════════════════════════════════
    # SLIDE 6: IC Scatter
    # ══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, WHITE)
    add_header(slide, "Indicator Testing: Information Coefficient (IC)", prs)
    add_image_fit(slide, "data/fig_ic_scatter.png",
                  Inches(0.3), Inches(CONTENT_TOP), Inches(12.7), Inches(5.0))
    add_caption(slide, "Negative slope confirms contrarian hypothesis  •  H1 SUPPORTED: IC significant at 24h (p < 0.05)")

    # ══════════════════════════════════════════════════════════════
    # SLIDE 7: Signal Timeline
    # ══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, WHITE)
    add_header(slide, "Signal Processing: Timeline (2021-2023)", prs)
    add_image_fit(slide, "data/fig_signal_timeline.png",
                  Inches(0.2), Inches(CONTENT_TOP), Inches(12.9), Inches(5.0))
    add_caption(slide, "Green = Long (z < -1.5)  •  Red = Short (z > 1.5)  •  Signals cluster around extreme events")

    # ══════════════════════════════════════════════════════════════
    # SLIDE 8: Equity Curves
    # ══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, WHITE)
    add_header(slide, "Trading Rules: Incremental Equity Curves", prs)
    add_image_fit(slide, "data/fig_equity_curves.png",
                  Inches(0.3), Inches(CONTENT_TOP), Inches(12.7), Inches(5.0))
    add_caption(slide, "Rule 0: Basic → Rule 1: +Max hold → Rule 2: +Stop loss → Rule 3: +OI filter  •  Gray: BTC buy-hold")

    # ══════════════════════════════════════════════════════════════
    # SLIDE 9: Grid Search
    # ══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, WHITE)
    add_header(slide, "Parameter Optimization: Grid Search (120 combinations)", prs)
    add_image_fit(slide, "data/fig_gridsearch_heatmap.png",
                  Inches(0.4), Inches(CONTENT_TOP), Inches(12.5), Inches(5.0))
    add_caption(slide, "Left: Threshold × Window heatmap  •  Right: Pardo check (best ≈ 2σ above mean)  •  H3 SUPPORTED")

    # ══════════════════════════════════════════════════════════════
    # SLIDE 10: Sensitivity
    # ══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, WHITE)
    add_header(slide, "Parameter Optimization: Sensitivity Analysis", prs)
    add_image_fit(slide, "data/fig_sensitivity.png",
                  Inches(0.3), Inches(CONTENT_TOP), Inches(12.7), Inches(5.0))
    add_caption(slide, "Flat = robust to choice  •  Sharp peak = overfit risk  •  Threshold is most sensitive")

    # ══════════════════════════════════════════════════════════════
    # SLIDE 11: Walk-Forward
    # ══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, WHITE)
    add_header(slide, "Walk-Forward Analysis (Out-of-Sample)", prs)
    add_image_fit(slide, "data/fig_walkforward.png",
                  Inches(0.3), Inches(CONTENT_TOP), Inches(12.7), Inches(5.0))
    add_caption(slide, "Most OOS windows positive  •  WF ratio 0.3-0.6  •  Parameter drift → regime dependence  •  H4 PARTIAL")

    # ══════════════════════════════════════════════════════════════
    # SLIDE 12: Bootstrap & Overfitting
    # ══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, WHITE)
    add_header(slide, "Overfitting Diagnostics", prs)
    add_image_fit(slide, "data/fig_bootstrap_sharpe.png",
                  Inches(0.4), Inches(CONTENT_TOP + 0.1), Inches(7.2), Inches(4.5))

    add_section_title(slide, "Deflated Sharpe Ratio", Inches(8), Inches(CONTENT_TOP))
    add_bullets(slide, [
        "Accounts for 120+ trials",
        "Adjusts for selection bias",
        "Some inflation expected"
    ], Inches(8), Inches(1.5), Inches(4.8), Inches(1.5), size=13)

    add_section_title(slide, "Top-N Removal", Inches(8), Inches(3.2))
    add_bullets(slide, [
        "Remove best 20 trades",
        "Still profitable",
        "Edge is distributed"
    ], Inches(8), Inches(3.6), Inches(4.8), Inches(1.5), size=13)

    add_section_title(slide, "PBO Cross-Test", Inches(8), Inches(5.2))
    add_text(slide, "Half-sample params show moderate stability",
             Inches(8), Inches(5.55), Inches(4.8), Inches(0.5), size=12, color=LIGHT_GRAY)

    # ══════════════════════════════════════════════════════════════
    # SLIDE 13: Extension
    # ══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, WHITE)
    add_header(slide, "Extension: Cross-Asset Validation", prs)
    add_image_fit(slide, "data/fig_cross_asset_corr.png",
                  Inches(0.3), Inches(CONTENT_TOP), Inches(5.5), Inches(4.0))
    add_image_fit(slide, "data/fig_extension_equity.png",
                  Inches(6.2), Inches(CONTENT_TOP), Inches(6.8), Inches(4.0))

    add_text(slide, "BTC-ETH funding correlation ~0.7+",
             Inches(0.5), Inches(5.2), Inches(5.5), Inches(0.3), size=12, color=LIGHT_GRAY)
    add_text(slide, "ETH positive Calmar with BTC params (no re-fit)",
             Inches(6.3), Inches(5.2), Inches(6.5), Inches(0.3), size=12, color=LIGHT_GRAY)

    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(MARGIN), Inches(5.7), Inches(12.5), Inches(0.65))
    box.fill.solid()
    box.fill.fore_color.rgb = UW_LIGHT_GOLD
    box.line.fill.background()
    add_text(slide, "Cross-asset validation: funding contrarian is a microstructure effect, not BTC-specific noise",
             Inches(0.6), Inches(5.85), Inches(12.2), Inches(0.4), size=14, bold=True, color=UW_PURPLE, align=PP_ALIGN.CENTER)

    # ══════════════════════════════════════════════════════════════
    # SLIDE 14: Conclusion
    # ══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, WHITE)
    add_header(slide, "Conclusion", prs)

    add_section_title(slide, "Key Findings", Inches(MARGIN), Inches(CONTENT_TOP))
    add_bullets(slide, [
        "Funding contrarian effect is real",
        "Economic basis: crowding + margin",
        "IC significant at 24h horizon",
        "OOS positive (40-60% degradation)",
        "Cross-asset validated on ETH"
    ], Inches(MARGIN), Inches(1.5), Inches(4), Inches(2.5), size=13)

    add_section_title(slide, "Backtest Results", Inches(4.6), Inches(CONTENT_TOP))
    add_bullets(slide, [
        "Calmar: ~0.11",
        "Sharpe: ~0.77",
        "Max DD: ~14%",
        "Win Rate: ~54%",
        "Trades: ~125 (4 years)"
    ], Inches(4.6), Inches(1.5), Inches(3.5), Inches(2.5), size=13)

    add_section_title(slide, "Hypothesis Results", Inches(8.5), Inches(CONTENT_TOP))
    add_bullets(slide, [
        "H1: ✓ Supported",
        "H2: ~ Mixed",
        "H3: ✓ Supported",
        "H4: ~ Partial"
    ], Inches(8.5), Inches(1.5), Inches(4), Inches(2), size=13)

    add_section_title(slide, "Risks & Limitations", Inches(MARGIN), Inches(4.2))
    add_bullets(slide, [
        "Tail events: LUNA-style collapse where funding stays extreme",
        "Regime dependence: better in high-volatility periods",
        "Market structure shift post-FTX may reduce effectiveness",
        "Crowding risk if strategy becomes widely adopted"
    ], Inches(MARGIN), Inches(4.6), Inches(12.5), Inches(1.8), size=13)

    # ══════════════════════════════════════════════════════════════
    # SLIDE 15: References
    # ══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, WHITE)
    add_header(slide, "References", prs)

    refs = [
        "Bailey, D. H., & López de Prado, M. (2014). The Deflated Sharpe Ratio: Correcting for",
        "    Selection Bias, Backtest Overfitting, and Non-Normality. Journal of Portfolio",
        "    Management, 40(5), 94-107.",
        "",
        "Brunnermeier, M. K., & Pedersen, L. H. (2009). Market Liquidity and Funding Liquidity.",
        "    The Review of Financial Studies, 22(6), 2201-2238.",
        "",
        "Liu, Y., Tsyvinski, A., & Wu, X. (2022). Common Risk Factors in Cryptocurrency.",
        "    The Journal of Finance, 77(2), 1133-1177.",
        "",
        "Pardo, R. (2008). The Evaluation and Optimization of Trading Strategies (2nd ed.).",
        "    John Wiley & Sons.",
        "",
        "Shleifer, A., & Vishny, R. W. (1997). The Limits of Arbitrage.",
        "    The Journal of Finance, 52(1), 35-55."
    ]

    box = slide.shapes.add_textbox(Inches(MARGIN), Inches(1.2), Inches(12.5), Inches(5.5))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(refs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_GRAY
        p.font.name = "Arial"

    # Save
    prs.save("CFRM522_Presentation.pptx")
    print(f"Saved: CFRM522_Presentation.pptx ({len(prs.slides)} slides)")

if __name__ == "__main__":
    create_presentation()
